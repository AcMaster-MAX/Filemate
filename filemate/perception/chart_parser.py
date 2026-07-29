"""ChartParser - 从文档中识别和提取图表信息。

支持识别以下类型的图表：
- 流程图
- 组织结构图
- 饼图、柱状图、折线图描述
- 时间线
- 思维导图

设计目标：
- 识别图表类型
- 提取图表标题和说明
- 提取图表中的关键数据点（如饼图的百分比、柱状图的值）
"""

from __future__ import annotations

import logging
import re
from dataclasses import dataclass, field
from enum import Enum, auto
from pathlib import Path
from typing import Any, Optional

logger = logging.getLogger(__name__)


class ChartType(Enum):
    """图表类型。"""
    UNKNOWN = auto()
    FLOWCHART = auto()        # 流程图
    ORGANIZATION = auto()     # 组织结构图
    PIE_CHART = auto()        # 饼图描述
    BAR_CHART = auto()        # 柱状图描述
    LINE_CHART = auto()       # 折线图描述
    TIMELINE = auto()         # 时间线
    MINDMAP = auto()          # 思维导图
    DIAGRAM = auto()          # 其他图示


@dataclass
class ChartDataPoint:
    """图表数据点。"""
    label: str           # 标签
    value: Optional[float] = None  # 数值（如果有）
    percentage: Optional[float] = None  # 百分比（如果有）
    description: str = ""  # 描述


@dataclass
class Chart:
    """提取的图表。"""
    chart_id: str
    chart_type: ChartType = ChartType.UNKNOWN
    title: str = ""           # 图表标题
    description: str = ""     # 图表说明
    data_points: list[ChartDataPoint] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)

    def to_task_elements(self) -> list[dict]:
        """转换为任务要素。"""
        elements = []

        if self.title:
            elements.append({"type": "chart_title", "text": self.title})

        if self.description:
            elements.append({"type": "chart_description", "text": self.description})

        for dp in self.data_points:
            element = {"type": "data_point", "label": dp.label}
            if dp.value is not None:
                element["value"] = dp.value
            if dp.percentage is not None:
                element["percentage"] = dp.percentage
            if dp.description:
                element["description"] = dp.description
            elements.append(element)

        return elements


class ChartParser:
    """图表解析器。

    用法::

        parser = ChartParser()
        charts = parser.extract_charts(file_path)
        for chart in charts:
            print(chart.title, chart.chart_type)
            for dp in chart.data_points:
                print(f"  - {dp.label}: {dp.percentage}")
    """

    def __init__(self) -> None:
        self._handlers = {
            "docx": self._parse_word,
            "doc": self._parse_word,
            "pptx": self._parse_ppt,
            "ppt": self._parse_ppt,
            "pdf": self._parse_pdf,
        }

        # 图表类型关键词
        self._chart_keywords = {
            ChartType.FLOWCHART: ["流程图", "流程", "步骤", "step", "process"],
            ChartType.ORGANIZATION: ["组织", "结构", "架构", "organization", "structure"],
            ChartType.TIMELINE: ["时间线", "日程", "timeline", "schedule", "进度"],
            ChartType.MINDMAP: ["思维导图", "脑图", "mind map", "思维"],
        }

        # 饼图数据模式
        self._pie_patterns = [
            re.compile(r"(\d+(?:\.\d+)?)\s*%"),
            re.compile(r"(\d+(?:\.\d+)?)\s*个百分点"),
        ]

        # 柱状图数据模式
        self._bar_patterns = [
            re.compile(r"(\d+(?:\.\d+)?)\s*[个台件条人元]"),
            re.compile(r"数量[:：]\s*(\d+)"),
        ]

    def extract_charts(self, file_path: str | Path) -> list[Chart]:
        """从文件中提取图表信息。

        Parameters
        ----------
        file_path : str | Path
            文件路径。

        Returns
        -------
        list[Chart]
            提取的图表列表。
        """
        path = Path(file_path)
        suffix = path.suffix.lstrip(".").lower()
        handler = self._handlers.get(suffix)

        if not handler:
            logger.warning("不支持提取图表的格式: .%s", suffix)
            return []

        try:
            return handler(path)
        except Exception as exc:
            logger.error("提取图表失败: %s", exc)
            return []

    def _parse_word(self, path: Path) -> list[Chart]:
        """从 Word 文档提取图表。"""
        charts = []
        try:
            from docx import Document

            doc = Document(str(path))

            # 查找形状（Word 中的图表可能以形状形式存在）
            for shape_idx, shape in enumerate(doc.inline_shapes):
                chart = self._extract_from_shape(shape, f"word_{shape_idx}")
                if chart:
                    charts.append(chart)

            # 从文本中推断图表信息
            text = "\n".join([p.text for p in doc.paragraphs])
            text_charts = self._infer_charts_from_text(text, f"word_text_{len(charts)}")
            charts.extend(text_charts)

            logger.info("从 Word 提取 %d 个图表: %s", len(charts), path.name)
        except ImportError:
            logger.warning("python-docx 未安装")
        except Exception as exc:
            logger.error("Word 图表解析失败: %s", exc)

        return charts

    def _parse_ppt(self, path: Path) -> list[Chart]:
        """从 PPT 提取图表。"""
        charts = []
        try:
            from pptx import Presentation

            prs = Presentation(str(path))

            for slide_idx, slide in enumerate(prs.slides):
                # 处理图表形状
                for shape_idx, shape in enumerate(slide.shapes):
                    # 检查是否为图表
                    if hasattr(shape, "chart"):
                        chart = self._extract_from_pptx_chart(
                            shape.chart,
                            slide_idx,
                            shape_idx,
                        )
                        if chart:
                            charts.append(chart)

                # 从文本框推断
                for tb_idx, tb in enumerate(slide.shapes):
                    if hasattr(tb, "text_frame"):
                        text = tb.text_frame.text
                        text_charts = self._infer_charts_from_text(
                            text,
                            f"ppt_text_{slide_idx}_{tb_idx}",
                        )
                        charts.extend(text_charts)

            logger.info("从 PPT 提取 %d 个图表: %s", len(charts), path.name)
        except ImportError:
            logger.warning("python-pptx 未安装")
        except Exception as exc:
            logger.error("PPT 图表解析失败: %s", exc)

        return charts

    def _parse_pdf(self, path: Path) -> list[Chart]:
        """从 PDF 提取图表。"""
        charts = []
        try:
            import pdfplumber

            with pdfplumber.open(str(path)) as pdf:
                for page_idx, page in enumerate(pdf.pages):
                    # 提取图片（图表可能以图片形式存在）
                    images = page.images
                    if images:
                        chart = Chart(
                            chart_id=f"pdf_img_{page_idx}",
                            chart_type=ChartType.DIAGRAM,
                            title=f"图片图表（第 {page_idx + 1} 页）",
                            metadata={"type": "image", "count": len(images)},
                        )
                        charts.append(chart)

                    # 从文本推断
                    text = page.extract_text()
                    if text:
                        text_charts = self._infer_charts_from_text(
                            text,
                            f"pdf_text_{page_idx}",
                        )
                        charts.extend(text_charts)

            logger.info("从 PDF 提取 %d 个图表: %s", len(charts), path.name)
        except ImportError:
            logger.warning("pdfplumber 未安装")
        except Exception as exc:
            logger.error("PDF 图表解析失败: %s", exc)

        return charts

    def _extract_from_shape(self, shape, chart_id: str) -> Optional[Chart]:
        """从 Word 形状提取图表。"""
        # Word 中的 inline_shapes 有 type 属性
        # 但需要具体处理不同的图表类型
        return None

    def _extract_from_pptx_chart(self, chart, slide_idx: int, shape_idx: int) -> Optional[Chart]:
        """从 PPTX 图表提取信息。"""
        try:
            # 获取图表标题
            title = ""
            if hasattr(chart, "title"):
                title = str(chart.title.text_frame.text) if chart.title else ""

            # 获取图表类型
            chart_type = self._detect_chart_type(title)

            # 提取数据点
            data_points = []
            if hasattr(chart, "series"):
                for series in chart.series:
                    for point in series.points:
                        label = ""
                        value = None

                        # 尝试获取标签
                        if hasattr(point, "category"):
                            label = str(point.category)

                        # 尝试获取值
                        if hasattr(point, "value"):
                            value = float(point.value) if point.value else None

                        data_points.append(ChartDataPoint(
                            label=label,
                            value=value,
                        ))

            return Chart(
                chart_id=f"chart_{slide_idx}_{shape_idx}",
                chart_type=chart_type,
                title=title,
                data_points=data_points,
                metadata={"source": "pptx"},
            )
        except Exception as exc:
            logger.debug("PPTX 图表提取失败: %s", exc)
            return None

    def _detect_chart_type(self, title: str) -> ChartType:
        """根据标题检测图表类型。"""
        title_lower = title.lower()

        for chart_type, keywords in self._chart_keywords.items():
            if any(kw in title_lower for kw in keywords):
                return chart_type

        return ChartType.UNKNOWN

    def _infer_charts_from_text(self, text: str, base_id: str) -> list[Chart]:
        """从文本中推断图表信息。"""
        charts = []
        lines = text.split("\n")

        # 检测列表结构（可能是饼图或流程图）
        in_list = False
        list_items = []
        list_title = ""

        for i, line in enumerate(lines):
            stripped = line.strip()
            if not stripped:
                continue

            # 检查是否是列表项（以数字或符号开头）
            if re.match(r"^[\d•‣◦⁃∙]\s+", stripped):
                if not in_list:
                    in_list = True
                    list_items = []
                    # 尝试获取前一行作为标题
                    if i > 0:
                        list_title = lines[i - 1].strip()

                list_items.append(stripped)
            else:
                if in_list and list_items:
                    # 分析列表，推断图表
                    chart = self._analyze_list_as_chart(
                        list_title,
                        list_items,
                        f"{base_id}_{i}",
                    )
                    if chart:
                        charts.append(chart)

                in_list = False
                list_items = []
                list_title = ""

        return charts

    def _analyze_list_as_chart(
        self,
        title: str,
        items: list[str],
        chart_id: str,
    ) -> Optional[Chart]:
        """将列表分析为图表。"""
        if not items:
            return None

        # 检测是否为饼图（包含百分比）
        percentages = []
        bar_values = []

        for item in items:
            # 匹配百分比
            for pattern in self._pie_patterns:
                match = pattern.search(item)
                if match:
                    try:
                        pct = float(match.group(1))
                        percentages.append(pct)
                    except ValueError:
                        pass

            # 匹配柱状图数值
            for pattern in self._bar_patterns:
                match = pattern.search(item)
                if match:
                    try:
                        val = float(match.group(1))
                        bar_values.append(val)
                    except ValueError:
                        pass

        if percentages:
            # 认为是饼图
            data_points = []
            for item in items:
                label = re.sub(r"^\d+[\.\)]\s*", "", item)
                # 提取百分比
                pct = 0
                for pattern in self._pie_patterns:
                    match = pattern.search(item)
                    if match:
                        try:
                            pct = float(match.group(1))
                        except ValueError:
                            pass
                        label = pattern.sub("", label)

                data_points.append(ChartDataPoint(
                    label=label.strip(),
                    percentage=pct,
                ))

            return Chart(
                chart_id=chart_id,
                chart_type=ChartType.PIE_CHART,
                title=title or "饼图",
                data_points=data_points,
                metadata={"inferred": True},
            )

        if bar_values:
            # 认为是柱状图
            data_points = []
            for i, item in enumerate(items):
                label = re.sub(r"^\d+[\.\)]\s*", "", item)
                value = bar_values[i] if i < len(bar_values) else None

                data_points.append(ChartDataPoint(
                    label=label.strip(),
                    value=value,
                ))

            return Chart(
                chart_id=chart_id,
                chart_type=ChartType.BAR_CHART,
                title=title or "柱状图",
                data_points=data_points,
                metadata={"inferred": True},
            )

        # 检测流程图（包含步骤词）
        step_keywords = ["第一步", "第二步", "首先", "然后", "最后", "step", "phase"]
        if any(any(kw in item.lower() for kw in step_keywords) for item in items):
            data_points = [
                ChartDataPoint(label=item.strip()) for item in items
            ]

            return Chart(
                chart_id=chart_id,
                chart_type=ChartType.FLOWCHART,
                title=title or "流程",
                data_points=data_points,
                metadata={"inferred": True},
            )

        # 检测时间线
        date_patterns = [
            r"\d{4}[-/年]\d{1,2}[-/月]\d{1,2}日?",
            r"\d{1,2}[-/月]\d{1,2}日?",
        ]
        has_dates = any(
            any(re.search(p, item) for p in date_patterns)
            for item in items
        )
        if has_dates or any(
            kw in title.lower()
            for kw in ["时间", "日程", "timeline", "schedule"]
        ):
            data_points = [
                ChartDataPoint(label=item.strip()) for item in items
            ]

            return Chart(
                chart_id=chart_id,
                chart_type=ChartType.TIMELINE,
                title=title or "时间线",
                data_points=data_points,
                metadata={"inferred": True},
            )

        return None

    def extract_task_charts(self, file_path: str | Path) -> list[dict[str, Any]]:
        """提取可能与任务相关的图表。"""
        charts = self.extract_charts(file_path)
        task_charts = []

        task_keywords = [
            "任务", "task", "截止", "deadline", "due",
            "提交", "submit", "完成", "finish",
            "计划", "plan", "进度", "progress",
        ]

        for chart in charts:
            # 检查标题
            if any(kw in chart.title.lower() for kw in task_keywords):
                task_charts.append({
                    "type": "title_match",
                    "chart": chart,
                })
                continue

            # 检查数据点
            for dp in chart.data_points:
                if any(kw in dp.label.lower() for kw in task_keywords):
                    task_charts.append({
                        "type": "data_match",
                        "chart": chart,
                    })
                    break

        return task_charts