"""TableReader - 从文档中提取表格数据。

支持从以下格式中提取表格：
- Word (.docx, .doc)
- PDF
- PPT

设计目标：
- 识别表格结构（行、列、单元格）
- 处理合并单元格
- 提取表格标题和说明
- 与任务提取流程集成
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Optional

logger = logging.getLogger(__name__)


@dataclass
class TableCell:
    """表格单元格。"""
    row: int
    col: int
    text: str
    is_header: bool = False
    is_merged: bool = False


@dataclass
class Table:
    """提取的表格。"""
    table_id: str
    caption: str = ""           # 表格标题
    rows: int = 0              # 行数
    cols: int = 0               # 列数
    cells: list[TableCell] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)

    def to_markdown(self) -> str:
        """转换为 Markdown 表格格式。"""
        if not self.cells:
            return ""

        # 构建表格结构
        grid = {}
        for cell in self.cells:
            key = (cell.row, cell.col)
            grid[key] = cell.text

        lines = []
        for r in range(self.rows):
            row_cells = []
            for c in range(self.cols):
                text = grid.get((r, c), "")
                row_cells.append(text)
            lines.append("| " + " | ".join(row_cells) + " |")

        # 添加分隔行
        if lines:
            lines.insert(1, "| " + " | ".join(["---"] * self.cols) + " |")

        # 添加标题
        if self.caption:
            lines.insert(0, f"**{self.caption}**\n")

        return "\n".join(lines)

    def extract_headers(self) -> list[str]:
        """提取表头。"""
        return [c.text for c in self.cells if c.is_header]

    def extract_data_rows(self) -> list[dict[str, str]]:
        """提取数据行（字典格式）。"""
        headers = self.extract_headers()
        if not headers:
            return []

        data = []
        for r in range(self.rows):
            row_data = {}
            for c in range(self.cols):
                cell_text = ""
                for cell in self.cells:
                    if cell.row == r and cell.col == c:
                        cell_text = cell.text
                        break
                if c < len(headers):
                    row_data[headers[c]] = cell_text
            if row_data:
                data.append(row_data)
        return data


class TableReader:
    """表格读取器。

    用法::

        reader = TableReader()
        tables = reader.extract_tables(file_path)
        for table in tables:
            print(table.to_markdown())
            for row in table.extract_data_rows():
                print(row)
    """

    def __init__(self) -> None:
        self._handlers = {
            "docx": self._read_word,
            "doc": self._read_word,
            "pdf": self._read_pdf,
            "pptx": self._read_ppt,
            "ppt": self._read_ppt,
        }

    def extract_tables(self, file_path: str | Path) -> list[Table]:
        """从文件中提取表格。

        Parameters
        ----------
        file_path : str | Path
            文件路径。

        Returns
        -------
        list[Table]
            提取的表格列表。
        """
        path = Path(file_path)
        suffix = path.suffix.lstrip(".").lower()
        handler = self._handlers.get(suffix)

        if not handler:
            logger.warning("不支持提取表格的格式: .%s", suffix)
            return []

        try:
            return handler(path)
        except Exception as exc:
            logger.error("提取表格失败: %s", exc)
            return []

    def _read_word(self, path: Path) -> list[Table]:
        """从 Word 文档提取表格。"""
        tables = []
        try:
            import docx
            from docx.table import Table as DocxTable

            doc = docx.Document(str(path))

            for idx, docx_table in enumerate(doc.tables):
                table = self._convert_docx_table(docx_table, idx)
                tables.append(table)

            logger.info("从 Word 提取 %d 个表格: %s", len(tables), path.name)
        except ImportError:
            logger.warning("python-docx 未安装，无法解析 Word 表格")
        except Exception as exc:
            logger.error("Word 表格解析失败: %s", exc)

        return tables

    def _convert_docx_table(self, docx_table: DocxTable, table_id: int) -> Table:
        """转换 Word 表格为通用格式。"""
        cells = []
        rows = len(docx_table.rows)
        cols = len(docx_table.columns) if docx_table.rows else 0

        # 检查第一行是否为表头（简单启发式）
        first_row_texts = [cell.text.strip() for cell in docx_table.rows[0].cells]

        for r_idx, row in enumerate(docx_table.rows):
            for c_idx, cell in enumerate(row.cells):
                text = cell.text.strip()
                is_header = (r_idx == 0 and any(first_row_texts))
                cells.append(TableCell(
                    row=r_idx,
                    col=c_idx,
                    text=text,
                    is_header=is_header,
                ))

        return Table(
            table_id=f"table_{table_id}",
            rows=rows,
            cols=cols,
            cells=cells,
        )

    def _read_pdf(self, path: Path) -> list[Table]:
        """从 PDF 提取表格。"""
        tables = []
        try:
            import pdfplumber

            with pdfplumber.open(str(path)) as pdf:
                for page_idx, page in enumerate(pdf.pages):
                    page_tables = page.extract_tables()
                    if page_tables:
                        for table_idx, table_data in enumerate(page_tables):
                            table = self._convert_pdf_table(
                                table_data,
                                page_idx,
                                table_idx,
                            )
                            if table.cells:
                                tables.append(table)

            logger.info("从 PDF 提取 %d 个表格: %s", len(tables), path.name)
        except ImportError:
            logger.warning("pdfplumber 未安装，无法解析 PDF 表格")
        except Exception as exc:
            logger.error("PDF 表格解析失败: %s", exc)

        return tables

    def _convert_pdf_table(self, table_data: list, page_idx: int, table_idx: int) -> Table:
        """转换 PDF 表格为通用格式。"""
        cells = []
        rows = len(table_data) if table_data else 0
        cols = len(table_data[0]) if table_data and table_data[0] else 0

        for r_idx, row in enumerate(table_data):
            for c_idx, cell_text in enumerate(row):
                text = str(cell_text).strip() if cell_text else ""
                is_header = (r_idx == 0)
                cells.append(TableCell(
                    row=r_idx,
                    col=c_idx,
                    text=text,
                    is_header=is_header,
                ))

        return Table(
            table_id=f"pdf_{page_idx}_{table_idx}",
            rows=rows,
            cols=cols,
            cells=cells,
            metadata={"source": "pdf", "page": page_idx},
        )

    def _read_ppt(self, path: Path) -> list[Table]:
        """从 PPT 提取表格。"""
        tables = []
        try:
            from pptx import Presentation

            prs = Presentation(str(path))

            for slide_idx, slide in enumerate(prs.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, "table"):
                        table = self._convert_pptx_table(
                            shape.table,
                            slide_idx,
                            shape_idx,
                        )
                        if table.cells:
                            tables.append(table)

            logger.info("从 PPT 提取 %d 个表格: %s", len(tables), path.name)
        except ImportError:
            logger.warning("python-pptx 未安装，无法解析 PPT 表格")
        except Exception as exc:
            logger.error("PPT 表格解析失败: %s", exc)

        return tables

    def _convert_pptx_table(self, ppt_table, slide_idx: int, shape_idx: int) -> Table:
        """转换 PPT 表格为通用格式。"""
        cells = []
        rows = len(ppt_table.rows)
        cols = len(ppt_table.columns) if ppt_table.rows else 0

        for r_idx, row in enumerate(ppt_table.rows):
            for c_idx, cell in enumerate(row.cells):
                text = cell.text.strip()
                is_header = (r_idx == 0)
                cells.append(TableCell(
                    row=r_idx,
                    col=c_idx,
                    text=text,
                    is_header=is_header,
                ))

        return Table(
            table_id=f"ppt_{slide_idx}_{shape_idx}",
            rows=rows,
            cols=cols,
            cells=cells,
            metadata={"source": "pptx", "slide": slide_idx},
        )

    def extract_task_tables(self, file_path: str | Path) -> list[dict[str, Any]]:
        """提取可能与任务相关的表格（如截止日期表、任务分配表）。"""
        tables = self.extract_tables(file_path)
        task_tables = []

        keywords = [
            "截止", "deadline", "due",
            "任务", "task", "assignment",
            "提交", "submit", "提交时间",
            "评分", "grade", "成绩",
            "名单", "list", "分组",
        ]

        for table in tables:
            # 检查表格标题
            if any(kw in table.caption.lower() for kw in keywords):
                task_tables.append({
                    "type": "caption_match",
                    "table": table,
                })
                continue

            # 检查表头
            headers = table.extract_headers()
            if any(any(kw in h.lower() for kw in keywords) for h in headers):
                task_tables.append({
                    "type": "header_match",
                    "table": table,
                })
                continue

            # 检查单元格内容
            for cell in table.cells:
                if any(kw in cell.text.lower() for kw in keywords):
                    task_tables.append({
                        "type": "content_match",
                        "table": table,
                    })
                    break

        return task_tables